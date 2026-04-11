from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptxcompose.composer import Composer


def slide_text(slide) -> str:
    values = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text and shape.text.strip():
            values.append(" ".join(shape.text.strip().split()))
    return " || ".join(values)


def find_insert_index(main: Presentation, after_title: str) -> int:
    key = after_title.lower()
    for index, slide in enumerate(main.slides):
        if key in slide_text(slide).lower():
            return index + 1

    for index, slide in enumerate(main.slides):
        if "stay tuned" in slide_text(slide).lower():
            return index

    return len(main.slides)


def clone_slides_from_eval(main_path: Path, eval_path: Path, output_path: Path, insert_index: int) -> None:
    main = Presentation(str(main_path))
    eval_prs = Presentation(str(eval_path))

    composer = Composer(main)
    composer.append(eval_prs)
    composer.save(str(output_path))

    merged = Presentation(str(output_path))
    sld_id_list = merged.slides._sldIdLst
    eval_count = len(eval_prs.slides)
    moved = [sld_id_list[-eval_count + i] for i in range(eval_count)]

    for _ in range(eval_count):
        sld_id_list.remove(sld_id_list[-1])

    for offset, slide_id in enumerate(moved):
        sld_id_list.insert(insert_index + offset, slide_id)

    merged.save(str(output_path))


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Merge exported Slidev eval PPTX into main PPTX")
    parser.add_argument("--main", required=True, type=Path, help="Path to main PPTX")
    parser.add_argument("--eval", required=True, type=Path, help="Path to eval PPTX exported from Slidev")
    parser.add_argument("--out", required=True, type=Path, help="Output merged PPTX path")
    parser.add_argument(
        "--after-title",
        default="Evaluation methodology",
        help="Insert eval deck after slide containing this title text",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()

    if not args.main.exists():
        raise FileNotFoundError(f"Main PPTX not found: {args.main}")
    if not args.eval.exists():
        raise FileNotFoundError(f"Eval PPTX not found: {args.eval}")

    main_prs = Presentation(str(args.main))
    insert_index = find_insert_index(main_prs, args.after_title)

    clone_slides_from_eval(args.main, args.eval, args.out, insert_index)
    print(f"Merged deck written to: {args.out}")


if __name__ == "__main__":
    main()
